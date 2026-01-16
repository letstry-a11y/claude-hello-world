
"""
HTML to PowerPoint Converter (通用版)

将 HTML 文件转换为 PowerPoint 演示文稿。
支持多种 HTML 结构和布局类型。

功能特点:
- 自动检测幻灯片分隔 (slide-container, section, article, hr)
- 支持多种内容类型 (标题, 段落, 列表, 图片, 表格)
- 可自定义主题颜色
- 提供 GUI 界面和命令行两种使用方式

Usage:
    # GUI 模式
    python html_to_pptx.py

    # 命令行模式
    python html_to_pptx.py input.html output.pptx
"""

import os
import re
import tempfile
import requests
from urllib.parse import urljoin, urlparse
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ============== 配置常量 ==============

# FontAwesome 图标映射
ICON_MAP = {
    "fa-check-circle": "✓", "fa-circle-check": "✓", "fa-check": "✓",
    "fa-exclamation-triangle": "⚠", "fa-triangle-exclamation": "⚠", "fa-warning": "⚠",
    "fa-globe-asia": "🌏", "fa-globe": "🌐", "fa-earth": "🌍",
    "fa-file-medical-alt": "📋", "fa-file-medical": "📋", "fa-file": "📄",
    "fa-robot": "🤖", "fa-microchip": "💡", "fa-lightbulb": "💡",
    "fa-star": "★", "fa-heart": "❤", "fa-arrow-right": "→",
    "fa-arrow-left": "←", "fa-info-circle": "ℹ", "fa-question-circle": "?",
    "fa-times": "✗", "fa-close": "✗", "fa-user": "👤", "fa-users": "👥",
    "fa-cog": "⚙", "fa-gear": "⚙", "fa-chart-bar": "📊", "fa-chart-line": "📈",
}

# 默认主题颜色
class ThemeColors:
    """可自定义的主题颜色"""
    def __init__(self, primary="#003366", accent="#0066CC",
                 text="#334155", muted="#64748B", success="#10B981", warning="#F59E0B"):
        self.primary = self._hex_to_rgb(primary)
        self.accent = self._hex_to_rgb(accent)
        self.text = self._hex_to_rgb(text)
        self.muted = self._hex_to_rgb(muted)
        self.success = self._hex_to_rgb(success)
        self.warning = self._hex_to_rgb(warning)
        self.white = RGBColor(0xFF, 0xFF, 0xFF)
        self.black = RGBColor(0x00, 0x00, 0x00)

    def _hex_to_rgb(self, hex_color):
        hex_color = hex_color.lstrip('#')
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )

# 默认主题
DEFAULT_THEME = ThemeColors()


# ============== 工具函数 ==============

def download_image(url, temp_dir, base_url=None):
    """下载图片到临时目录"""
    try:
        # 处理相对路径
        if base_url and not url.startswith(('http://', 'https://', 'data:')):
            url = urljoin(base_url, url)

        # 跳过 data URL
        if url.startswith('data:'):
            return None

        response = requests.get(url, timeout=30, headers={
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
        })
        response.raise_for_status()

        # 生成文件名
        parsed = urlparse(url)
        filename = os.path.basename(parsed.path) or f"image_{hash(url) % 10000}.jpg"
        if '.' not in filename:
            filename += '.jpg'

        filepath = os.path.join(temp_dir, filename)
        with open(filepath, 'wb') as f:
            f.write(response.content)
        return filepath
    except Exception as e:
        print(f"警告: 无法下载图片 {url}: {e}")
        return None


def get_icon_unicode(element):
    """从元素中提取 FontAwesome 图标并转换为 Unicode"""
    icon_elem = element.find('i', class_=re.compile(r'fa'))
    if icon_elem:
        for cls in icon_elem.get('class', []):
            if cls.startswith('fa-') and cls in ICON_MAP:
                return ICON_MAP[cls]
    return "•"


def clean_text(text):
    """清理文本内容"""
    if not text:
        return ""
    # 移除多余空白
    text = re.sub(r'\s+', ' ', text).strip()
    return text


def extract_color_from_style(style_str):
    """从 style 属性中提取颜色"""
    if not style_str:
        return None
    match = re.search(r'color:\s*([#\w]+)', style_str)
    if match:
        color = match.group(1)
        if color.startswith('#') and len(color) == 7:
            return RGBColor(
                int(color[1:3], 16),
                int(color[3:5], 16),
                int(color[5:7], 16)
            )
    return None


# ============== 幻灯片创建函数 ==============

def add_text_box(slide, left, top, width, height, text,
                 font_size=14, color=None, bold=False, align=PP_ALIGN.LEFT,
                 theme=DEFAULT_THEME):
    """添加文本框"""
    if color is None:
        color = theme.text

    shape = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = shape.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = clean_text(text)
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align

    return shape


def add_title_subtitle(slide, title, subtitle, theme=DEFAULT_THEME):
    """添加标题和副标题"""
    # 标题
    if title:
        add_text_box(slide, 0.8, 0.5, 11.7, 0.7, title,
                     font_size=28, color=theme.primary, bold=True, theme=theme)

    # 副标题
    if subtitle:
        add_text_box(slide, 0.8, 1.15, 11.7, 0.5, subtitle,
                     font_size=18, color=theme.muted, theme=theme)

        # 添加装饰线
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8), Inches(1.7), Inches(11.7), Inches(0.03)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = theme.accent
        line.line.fill.background()


def add_bullet_list(slide, items, start_y, theme=DEFAULT_THEME):
    """添加项目符号列表"""
    y_pos = start_y

    for item in items:
        icon = item.get('icon', '•')
        title = item.get('title', '')
        text = item.get('text', '')
        icon_color = item.get('icon_color', theme.accent)

        # 图标
        add_text_box(slide, 0.8, y_pos, 0.4, 0.4, icon,
                     font_size=18, color=icon_color, bold=True, theme=theme)

        # 标题
        if title:
            add_text_box(slide, 1.3, y_pos, 10.5, 0.35, title,
                         font_size=16, color=theme.primary, bold=True, theme=theme)
            y_pos += 0.35

        # 内容
        if text:
            add_text_box(slide, 1.3, y_pos, 10.5, 0.8, text,
                         font_size=14, color=theme.text, theme=theme)

        y_pos += 1.0 if not title else 0.8

    return y_pos


def add_image(slide, img_path, left, top, width=None, height=None):
    """添加图片"""
    if img_path and os.path.exists(img_path):
        try:
            if width and height:
                slide.shapes.add_picture(img_path, Inches(left), Inches(top),
                                        width=Inches(width), height=Inches(height))
            elif width:
                slide.shapes.add_picture(img_path, Inches(left), Inches(top),
                                        width=Inches(width))
            else:
                slide.shapes.add_picture(img_path, Inches(left), Inches(top))
            return True
        except Exception as e:
            print(f"警告: 无法添加图片: {e}")
    return False


def add_card(slide, x, y, width, height, title, text, icon="•", theme=DEFAULT_THEME):
    """添加卡片样式内容"""
    # 左侧装饰条
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y + 0.3), Inches(0.05), Inches(height - 0.6)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme.accent
    bar.line.fill.background()

    # 图标框
    icon_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x + 0.15), Inches(y + 0.15), Inches(0.5), Inches(0.5)
    )
    icon_box.fill.solid()
    icon_box.fill.fore_color.rgb = theme.primary
    icon_box.line.fill.background()

    # 图标文字
    add_text_box(slide, x + 0.15, y + 0.22, 0.5, 0.4, icon,
                 font_size=16, color=theme.white, align=PP_ALIGN.CENTER, theme=theme)

    # 标题
    add_text_box(slide, x + 0.75, y + 0.2, width - 0.9, 0.4, title,
                 font_size=15, color=theme.primary, bold=True, theme=theme)

    # 内容
    add_text_box(slide, x + 0.15, y + 0.7, width - 0.3, height - 0.9, text,
                 font_size=12, color=theme.muted, theme=theme)


def add_footer(slide, left_text, right_text, theme=DEFAULT_THEME):
    """添加页脚"""
    if left_text:
        add_text_box(slide, 0.8, 7.0, 2.0, 0.3, left_text,
                     font_size=10, color=theme.muted, theme=theme)

    if right_text:
        add_text_box(slide, 9.5, 7.0, 3.0, 0.3, right_text,
                     font_size=10, color=theme.muted, bold=True,
                     align=PP_ALIGN.RIGHT, theme=theme)


# ============== HTML 解析函数 ==============

def find_slides(soup):
    """查找 HTML 中的幻灯片分隔"""
    # 尝试不同的幻灯片容器
    containers = (
        soup.find_all('div', class_='slide-container') or
        soup.find_all('div', class_='slide') or
        soup.find_all('section') or
        soup.find_all('article')
    )

    if containers:
        return containers

    # 如果没有明确的容器，尝试用 hr 分隔
    body = soup.find('body') or soup

    # 检查是否有 hr 分隔
    hrs = body.find_all('hr')
    if hrs:
        slides = []
        current_content = []
        for child in body.children:
            if child.name == 'hr':
                if current_content:
                    wrapper = soup.new_tag('div')
                    for c in current_content:
                        wrapper.append(c.extract() if hasattr(c, 'extract') else c)
                    slides.append(wrapper)
                    current_content = []
            else:
                current_content.append(child)
        if current_content:
            wrapper = soup.new_tag('div')
            for c in current_content:
                if hasattr(c, 'extract'):
                    wrapper.append(c)
            slides.append(wrapper)
        return slides if slides else [body]

    # 没有分隔，整个 body 作为一个幻灯片
    return [body]


def extract_slide_content(container):
    """从容器中提取幻灯片内容"""
    content = {
        'title': '',
        'subtitle': '',
        'items': [],
        'images': [],
        'cards': [],
        'tables': [],
        'footer_left': '',
        'footer_right': '',
        'layout': 'auto'
    }

    # 提取标题
    title_elem = (
        container.find(class_='slide-title') or
        container.find('h1') or
        container.find('h2')
    )
    if title_elem:
        content['title'] = clean_text(title_elem.get_text())

    # 提取副标题
    subtitle_elem = (
        container.find(class_='slide-subtitle') or
        container.find('h3')
    )
    if subtitle_elem and subtitle_elem != title_elem:
        content['subtitle'] = clean_text(subtitle_elem.get_text())

    # 提取列表项
    for ul in container.find_all(['ul', 'ol']):
        for li in ul.find_all('li', recursive=False):
            item = {'icon': '•', 'title': '', 'text': '', 'icon_color': None}

            # 检查类名以确定图标类型
            li_classes = li.get('class', [])
            if 'strength' in li_classes or 'success' in li_classes:
                item['icon'] = '✓'
                item['icon_color'] = 'success'
            elif 'gap' in li_classes or 'warning' in li_classes:
                item['icon'] = '⚠'
                item['icon_color'] = 'warning'
            else:
                item['icon'] = get_icon_unicode(li)

            # 提取标题和文本
            strong = li.find(['strong', 'b'])
            if strong:
                item['title'] = clean_text(strong.get_text())
                # 获取 strong 标签之后的所有文本内容
                remaining_text = []
                for sibling in strong.next_siblings:
                    if hasattr(sibling, 'get_text'):
                        remaining_text.append(sibling.get_text())
                    elif isinstance(sibling, str):
                        remaining_text.append(sibling)
                item['text'] = clean_text(' '.join(remaining_text))
            else:
                item['text'] = clean_text(li.get_text())

            content['items'].append(item)

    # 提取图片
    for img in container.find_all('img'):
        src = img.get('src', '')
        alt = img.get('alt', '')
        if src:
            content['images'].append({'src': src, 'alt': alt})

    # 提取卡片 (tile-card, roadmap-card, card)
    # 使用更精确的匹配，只匹配顶级卡片类，避免匹配子元素
    card_selectors = [
        ('div', 'tile-card'),
        ('div', 'roadmap-card'),
        ('div', 'card'),
    ]
    seen_cards = set()
    for tag, cls in card_selectors:
        for card in container.find_all(tag, class_=cls):
            # 避免重复处理同一个卡片元素
            card_id = id(card)
            if card_id in seen_cards:
                continue
            seen_cards.add(card_id)

            card_data = {
                'title': '',
                'text': '',
                'icon': '•',
                'image': None
            }

            # 卡片标题
            card_title = card.find(['h3', 'h4', 'h5'])
            if card_title:
                card_data['title'] = clean_text(card_title.get_text())

            # 卡片内容
            card_text = card.find('p')
            if card_text:
                card_data['text'] = clean_text(card_text.get_text())

            # 卡片图标
            card_data['icon'] = get_icon_unicode(card)

            # 卡片图片
            card_img = card.find('img')
            if card_img and card_img.get('src'):
                card_data['image'] = card_img['src']

            content['cards'].append(card_data)

    # 提取表格
    for table in container.find_all('table'):
        table_data = []
        for row in table.find_all('tr'):
            row_data = []
            for cell in row.find_all(['th', 'td']):
                row_data.append(clean_text(cell.get_text()))
            if row_data:
                table_data.append(row_data)
        if table_data:
            content['tables'].append(table_data)

    # 提取页脚
    page_indicator = container.find(class_='page-indicator')
    footer = container.find(class_='footer')

    if page_indicator:
        content['footer_left'] = clean_text(page_indicator.get_text())
    if footer:
        content['footer_right'] = clean_text(footer.get_text())

    # 检测布局类型
    if container.find(class_='two-column'):
        content['layout'] = 'two-column'
    elif container.find(class_='tile-grid'):
        content['layout'] = 'tile-grid'
    elif container.find(class_='roadmap-grid'):
        content['layout'] = 'roadmap-grid'
    elif content['cards']:
        content['layout'] = 'cards'
    elif content['images'] and content['items']:
        content['layout'] = 'two-column'

    return content


# ============== 主转换函数 ==============

def create_slide(prs, content, temp_dir, base_url=None, theme=DEFAULT_THEME):
    """创建单个幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

    # 添加标题和副标题
    add_title_subtitle(slide, content['title'], content['subtitle'], theme)

    layout = content['layout']
    start_y = 2.0 if content['title'] else 0.8

    # 处理颜色映射
    for item in content['items']:
        if item['icon_color'] == 'success':
            item['icon_color'] = theme.success
        elif item['icon_color'] == 'warning':
            item['icon_color'] = theme.warning
        else:
            item['icon_color'] = theme.accent

    # 根据布局类型创建内容
    if layout == 'two-column':
        # 左侧列表
        if content['items']:
            add_bullet_list(slide, content['items'], start_y, theme)

        # 右侧图片
        if content['images']:
            img_src = content['images'][0]['src']
            img_path = download_image(img_src, temp_dir, base_url)
            add_image(slide, img_path, 7.3, start_y, 5.2, 4.0)

    elif layout in ('tile-grid', 'cards') and content['cards']:
        # 卡片网格布局
        cards = content['cards']

        if len(cards) <= 2:
            # 2 列布局
            positions = [(0.8, start_y), (7.0, start_y)]
            card_width, card_height = 5.5, 4.5
        else:
            # 2x2 布局
            positions = [
                (0.8, start_y), (6.9, start_y),
                (0.8, start_y + 2.4), (6.9, start_y + 2.4)
            ]
            card_width, card_height = 5.8, 2.2

        for i, card in enumerate(cards[:4]):
            if i >= len(positions):
                break
            x, y = positions[i]

            # 如果卡片有图片，先添加图片
            if card['image']:
                img_path = download_image(card['image'], temp_dir, base_url)
                if img_path:
                    add_image(slide, img_path, x, y, card_width, 2.5)
                    # 在图片下方添加标题和文字
                    add_text_box(slide, x, y + 2.6, card_width, 0.4, card['title'],
                                font_size=16, color=theme.primary, bold=True,
                                align=PP_ALIGN.CENTER, theme=theme)
                    add_text_box(slide, x, y + 3.0, card_width, 1.0, card['text'],
                                font_size=12, color=theme.muted,
                                align=PP_ALIGN.CENTER, theme=theme)
            else:
                # 无图片的卡片样式
                add_card(slide, x, y, card_width, card_height,
                        card['title'], card['text'], card['icon'], theme)

    elif layout == 'roadmap-grid' and content['cards']:
        # 路线图网格
        cards = content['cards']
        positions = [
            (0.8, start_y), (6.9, start_y),
            (0.8, start_y + 2.4), (6.9, start_y + 2.4)
        ]

        for i, card in enumerate(cards[:4]):
            if i >= len(positions):
                break
            x, y = positions[i]
            add_card(slide, x, y, 5.8, 2.2, card['title'], card['text'], card['icon'], theme)

    else:
        # 默认布局 - 列表 + 图片
        y_pos = start_y

        if content['items']:
            y_pos = add_bullet_list(slide, content['items'], y_pos, theme)

        # 添加图片
        for i, img_data in enumerate(content['images'][:2]):
            img_path = download_image(img_data['src'], temp_dir, base_url)
            if img_path:
                add_image(slide, img_path, 0.8 + i * 6.2, y_pos, 5.5, 3.0)

    # 添加页脚
    add_footer(slide, content['footer_left'], content['footer_right'], theme)

    return slide


def convert_html_to_pptx(html_path, output_path, theme=None, progress_callback=None):
    """
    将 HTML 转换为 PowerPoint

    Args:
        html_path: HTML 文件路径或 HTML 内容字符串
        output_path: 输出 PPTX 文件路径
        theme: 自定义主题颜色 (ThemeColors 实例)
        progress_callback: 进度回调函数 (current, total, message)

    Returns:
        输出文件路径
    """
    if theme is None:
        theme = DEFAULT_THEME

    # 读取 HTML
    base_url = None
    if os.path.isfile(html_path):
        base_url = f"file:///{os.path.dirname(os.path.abspath(html_path))}/"
        with open(html_path, 'r', encoding='utf-8') as f:
            html_content = f.read()
    else:
        html_content = html_path

    soup = BeautifulSoup(html_content, 'html.parser')

    # 创建演示文稿 (16:9 宽屏)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 查找幻灯片
    slide_containers = find_slides(soup)
    total_slides = len(slide_containers)

    if total_slides == 0:
        raise ValueError("HTML 中没有找到可转换的内容")

    if progress_callback:
        progress_callback(0, total_slides, "开始转换...")

    # 创建临时目录存放下载的图片
    with tempfile.TemporaryDirectory() as temp_dir:
        for i, container in enumerate(slide_containers):
            if progress_callback:
                progress_callback(i, total_slides, f"处理第 {i+1}/{total_slides} 页...")

            content = extract_slide_content(container)
            create_slide(prs, content, temp_dir, base_url, theme)

    # 保存文件
    prs.save(output_path)

    if progress_callback:
        progress_callback(total_slides, total_slides, "转换完成!")

    return output_path


# ============== GUI 界面 ==============

def create_gui():
    """创建图形用户界面"""
    import tkinter as tk
    from tkinter import ttk, filedialog, messagebox
    import threading

    class ConverterApp:
        def __init__(self, root):
            self.root = root
            self.root.title("HTML 转 PowerPoint 转换器")
            self.root.geometry("600x450")
            self.root.resizable(True, True)

            # 设置样式
            style = ttk.Style()
            style.configure('Title.TLabel', font=('Microsoft YaHei', 16, 'bold'))
            style.configure('Info.TLabel', font=('Microsoft YaHei', 9))

            self.setup_ui()

        def setup_ui(self):
            # 主框架
            main_frame = ttk.Frame(self.root, padding="20")
            main_frame.pack(fill=tk.BOTH, expand=True)

            # 标题
            title_label = ttk.Label(main_frame, text="HTML → PowerPoint 转换器",
                                   style='Title.TLabel')
            title_label.pack(pady=(0, 20))

            # 输入文件选择
            input_frame = ttk.LabelFrame(main_frame, text="输入文件", padding="10")
            input_frame.pack(fill=tk.X, pady=(0, 10))

            self.input_var = tk.StringVar()
            input_entry = ttk.Entry(input_frame, textvariable=self.input_var, width=50)
            input_entry.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(0, 10))

            input_btn = ttk.Button(input_frame, text="浏览...", command=self.browse_input)
            input_btn.pack(side=tk.RIGHT)

            # 输出文件选择
            output_frame = ttk.LabelFrame(main_frame, text="输出文件", padding="10")
            output_frame.pack(fill=tk.X, pady=(0, 10))

            self.output_var = tk.StringVar()
            output_entry = ttk.Entry(output_frame, textvariable=self.output_var, width=50)
            output_entry.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(0, 10))

            output_btn = ttk.Button(output_frame, text="浏览...", command=self.browse_output)
            output_btn.pack(side=tk.RIGHT)

            # 主题颜色设置
            theme_frame = ttk.LabelFrame(main_frame, text="主题颜色 (可选)", padding="10")
            theme_frame.pack(fill=tk.X, pady=(0, 10))

            color_frame = ttk.Frame(theme_frame)
            color_frame.pack(fill=tk.X)

            ttk.Label(color_frame, text="主色:").pack(side=tk.LEFT)
            self.primary_var = tk.StringVar(value="#003366")
            ttk.Entry(color_frame, textvariable=self.primary_var, width=10).pack(side=tk.LEFT, padx=(5, 20))

            ttk.Label(color_frame, text="强调色:").pack(side=tk.LEFT)
            self.accent_var = tk.StringVar(value="#0066CC")
            ttk.Entry(color_frame, textvariable=self.accent_var, width=10).pack(side=tk.LEFT, padx=5)

            # 进度条
            progress_frame = ttk.Frame(main_frame)
            progress_frame.pack(fill=tk.X, pady=(10, 10))

            self.progress_var = tk.DoubleVar()
            self.progress_bar = ttk.Progressbar(progress_frame, variable=self.progress_var,
                                                maximum=100, mode='determinate')
            self.progress_bar.pack(fill=tk.X)

            self.status_var = tk.StringVar(value="准备就绪")
            status_label = ttk.Label(progress_frame, textvariable=self.status_var,
                                    style='Info.TLabel')
            status_label.pack(pady=(5, 0))

            # 转换按钮
            btn_frame = ttk.Frame(main_frame)
            btn_frame.pack(pady=20)

            self.convert_btn = ttk.Button(btn_frame, text="开始转换",
                                         command=self.start_conversion)
            self.convert_btn.pack(side=tk.LEFT, padx=10)

            open_btn = ttk.Button(btn_frame, text="打开输出文件夹",
                                 command=self.open_output_folder)
            open_btn.pack(side=tk.LEFT, padx=10)

            # 说明信息
            info_text = """支持的 HTML 结构:
• 使用 <div class="slide-container"> 或 <section> 分隔幻灯片
• 使用 <h1>/<h2> 作为标题, <h3> 作为副标题
• 使用 <ul>/<ol> 创建列表, <img> 添加图片
• 使用 <hr> 分隔不同幻灯片"""

            info_label = ttk.Label(main_frame, text=info_text, style='Info.TLabel',
                                  justify=tk.LEFT)
            info_label.pack(pady=(10, 0), anchor=tk.W)

        def browse_input(self):
            filepath = filedialog.askopenfilename(
                title="选择 HTML 文件",
                filetypes=[("HTML 文件", "*.html *.htm"), ("所有文件", "*.*")]
            )
            if filepath:
                self.input_var.set(filepath)
                # 自动设置输出路径
                output_path = os.path.splitext(filepath)[0] + ".pptx"
                self.output_var.set(output_path)

        def browse_output(self):
            filepath = filedialog.asksaveasfilename(
                title="保存 PowerPoint 文件",
                defaultextension=".pptx",
                filetypes=[("PowerPoint 文件", "*.pptx"), ("所有文件", "*.*")]
            )
            if filepath:
                self.output_var.set(filepath)

        def update_progress(self, current, total, message):
            if total > 0:
                self.progress_var.set((current / total) * 100)
            self.status_var.set(message)
            self.root.update_idletasks()

        def start_conversion(self):
            input_path = self.input_var.get()
            output_path = self.output_var.get()

            if not input_path:
                messagebox.showerror("错误", "请选择输入的 HTML 文件")
                return

            if not output_path:
                messagebox.showerror("错误", "请指定输出的 PPTX 文件路径")
                return

            if not os.path.exists(input_path):
                messagebox.showerror("错误", f"输入文件不存在: {input_path}")
                return

            # 创建自定义主题
            try:
                theme = ThemeColors(
                    primary=self.primary_var.get(),
                    accent=self.accent_var.get()
                )
            except:
                theme = DEFAULT_THEME

            self.convert_btn.config(state='disabled')
            self.progress_var.set(0)

            def convert_thread():
                try:
                    convert_html_to_pptx(
                        input_path, output_path, theme,
                        progress_callback=self.update_progress
                    )
                    self.root.after(0, lambda: messagebox.showinfo(
                        "成功", f"转换完成!\n\n输出文件: {output_path}"
                    ))
                except Exception as e:
                    self.root.after(0, lambda: messagebox.showerror(
                        "转换失败", f"转换过程中出错:\n{str(e)}"
                    ))
                finally:
                    self.root.after(0, lambda: self.convert_btn.config(state='normal'))

            threading.Thread(target=convert_thread, daemon=True).start()

        def open_output_folder(self):
            output_path = self.output_var.get()
            if output_path:
                folder = os.path.dirname(output_path) or "."
                os.startfile(folder)

    root = tk.Tk()
    app = ConverterApp(root)
    root.mainloop()


# ============== 入口点 ==============

if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        # 无参数时启动 GUI
        print("启动图形界面...")
        create_gui()
    else:
        # 命令行模式
        input_file = sys.argv[1]
        output_file = sys.argv[2] if len(sys.argv) > 2 else input_file.replace('.html', '.pptx')

        print(f"输入: {input_file}")
        print(f"输出: {output_file}")

        def progress(current, total, msg):
            print(f"  [{current}/{total}] {msg}")

        try:
            convert_html_to_pptx(input_file, output_file, progress_callback=progress)
            print(f"\n转换成功! 输出文件: {output_file}")
        except Exception as e:
            print(f"\n转换失败: {e}")
            sys.exit(1)
